"""
Office / PDF 文档解析器。

支持 PDF / Word (.docx) / Excel (.xlsx) / PPT (.pptx) 格式。
使用 pypdf、python-docx、openpyxl、python-pptx 库。
"""
from __future__ import annotations

import logging
import uuid
from pathlib import Path
from typing import Any

from app.services.parser.base import Block, BlockType, ParsedDocument

logger = logging.getLogger(__name__)

# 支持的 MIME 类型
_OFFICE_MIME_TYPES = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}

# 扩展名 → MIME 映射
_EXT_TO_MIME = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class OfficeParser:
    """
    Office / PDF 文档解析器。

    实现 Parser 协议，支持 PDF / Word / Excel / PPT。
    """

    supported_types: set[str] = set(_OFFICE_MIME_TYPES.keys())

    async def parse(
        self,
        file_path: str,
        document_id: str,
        metadata: dict[str, Any],
    ) -> ParsedDocument:
        """解析文档文件，返回结构化内容。"""
        path = Path(file_path)
        ext = path.suffix.lower()

        logger.info("office_parser_start", extra={"file": str(path), "ext": ext})

        if ext == ".pdf":
            blocks, total_pages = self._parse_pdf(path)
        elif ext == ".docx":
            blocks, total_pages = self._parse_docx(path)
        elif ext == ".xlsx":
            blocks, total_pages = self._parse_xlsx(path)
        elif ext == ".pptx":
            blocks, total_pages = self._parse_pptx(path)
        else:
            blocks = [self._make_block(f"不支持的文件格式: {ext}", BlockType.PARAGRAPH, 0)]
            total_pages = 1

        title = path.stem
        source_name = path.name
        pages = list(range(1, total_pages + 1))
        tables = [b for b in blocks if b.block_type == BlockType.TABLE]

        logger.info(
            "office_parser_done",
            extra={"file": str(path), "block_count": len(blocks)},
        )

        return ParsedDocument(
            document_id=document_id,
            source_name=source_name,
            title=title,
            blocks=blocks,
            pages=pages,
            tables=tables,
            images=[],
            metadata=metadata,
            parser_used="office",
            total_pages=total_pages,
        )

    def _make_block(
        self,
        text: str,
        block_type: BlockType,
        order: int,
        page: int = 1,
        heading_path: str = "",
        metadata: dict[str, Any] | None = None,
    ) -> Block:
        return Block(
            block_id=f"blk_{uuid.uuid4().hex[:8]}",
            block_type=block_type,
            text=text,
            page_number=page,
            heading_path=heading_path,
            order_index=order,
            metadata=metadata or {},
        )

    def _parse_pdf(self, path: Path) -> tuple[list[Block], int]:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        blocks: list[Block] = []
        order = 0

        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                blocks.append(self._make_block(
                    text.strip(), BlockType.PARAGRAPH, order, page=i + 1,
                ))
                order += 1

        return blocks, len(reader.pages)

    def _parse_docx(self, path: Path) -> tuple[list[Block], int]:
        from docx import Document

        doc = Document(str(path))
        blocks: list[Block] = []
        order = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            block_type = BlockType.PARAGRAPH
            heading_path = ""
            if para.style and para.style.name:
                style_name = para.style.name.lower()
                if "heading" in style_name:
                    block_type = BlockType.HEADING
                    heading_path = text

            blocks.append(self._make_block(text, block_type, order, heading_path=heading_path))
            order += 1

        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                blocks.append(self._make_block("\n".join(rows), BlockType.TABLE, order))
                order += 1

        return blocks, 1

    def _parse_xlsx(self, path: Path) -> tuple[list[Block], int]:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        blocks: list[Block] = []
        order = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))

            if rows:
                blocks.append(self._make_block(
                    f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
                    BlockType.TABLE, order,
                    metadata={"sheet": sheet_name},
                ))
                order += 1

        wb.close()
        return blocks, 1

    def _parse_pptx(self, path: Path) -> tuple[list[Block], int]:
        from pptx import Presentation

        prs = Presentation(str(path))
        blocks: list[Block] = []
        order = 0

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))

            if texts:
                blocks.append(self._make_block(
                    "\n".join(texts), BlockType.PARAGRAPH, order, page=i + 1,
                ))
                order += 1

        return blocks, len(prs.slides)
